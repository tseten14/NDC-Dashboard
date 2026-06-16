#!/usr/bin/env python3
"""Generate June 17 NDC Data Explorer demo PowerPoint (diagrams + Climate TRACE)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "demo" / "deck.pptx"
SCREENSHOTS = ROOT / "docs" / "demo" / "screenshots"
OUT.parent.mkdir(parents=True, exist_ok=True)

LIVE_URL = "https://ndc-data-explorer-e051f914.vercel.app"
CLIMATE_TRACE_API = "https://api.climatetrace.org/v7"

# Theme (app greens)
DARK = RGBColor(30, 42, 36)
ACCENT = RGBColor(57, 107, 79)
ACCENT_LIGHT = RGBColor(72, 130, 95)
LIGHT = RGBColor(247, 249, 248)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(95, 110, 100)
PALE = RGBColor(220, 228, 224)
SUBTITLE = RGBColor(180, 200, 190)
OCEAN = RGBColor(46, 120, 160)
AMBER = RGBColor(180, 130, 40)

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

MARGIN_L = Inches(0.5)
MARGIN_R = Inches(0.5)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R

HEADER_H = Inches(0.72)
FOOTER_TOP = Inches(5.18)
FOOTER_H = Inches(0.405)
BODY_TOP = Inches(0.82)
BODY_BOTTOM = FOOTER_TOP - Inches(0.08)


def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_rect(slide, left, top, width, height, fill, line=None, *, radius=None):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def styled_textbox(slide, left, top, width, height, text, *, size=12, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.line_spacing = 1.15
    return box


def add_header(slide, title, subtitle=None):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, HEADER_H, DARK)
    styled_textbox(
        slide, MARGIN_L, Inches(0.14), CONTENT_W, Inches(0.34),
        title, size=20, bold=True, color=WHITE,
    )
    if subtitle:
        styled_textbox(
            slide, MARGIN_L, Inches(0.44), CONTENT_W, Inches(0.22),
            subtitle, size=10, color=SUBTITLE,
        )


def add_footer(slide, num):
    add_rect(slide, Inches(0), FOOTER_TOP, SLIDE_W, FOOTER_H, DARK)
    styled_textbox(
        slide, MARGIN_L, FOOTER_TOP + Inches(0.06), Inches(7.5), Inches(0.28),
        "Uganda NDC Data Explorer  ·  Demo 17 June 2026  ·  Data: Climate TRACE API v7",
        size=8, color=SUBTITLE,
    )
    styled_textbox(
        slide, Inches(8.85), FOOTER_TOP + Inches(0.06), Inches(0.65), Inches(0.28),
        f"{num:02d}", size=9, bold=True, color=WHITE, align=PP_ALIGN.RIGHT,
    )


def content_slide(prs, title, subtitle, num):
    slide = blank_slide(prs)
    add_header(slide, title, subtitle)
    add_footer(slide, num)
    return slide


def col_widths(n: int, gap=Inches(0.22)):
    total_gap = gap * (n - 1)
    w = (CONTENT_W - total_gap) / n
    lefts = [MARGIN_L + i * (w + gap) for i in range(n)]
    return lefts, w, gap


def add_three_cards(slide, cards, top=BODY_TOP, card_h=None):
    lefts, w, _ = col_widths(3)
    if card_h is None:
        card_h = BODY_BOTTOM - top
    for left, (head, body) in zip(lefts, cards):
        add_rect(slide, left, top, w, card_h, WHITE, PALE)
        add_rect(slide, left, top, w, Inches(0.38), ACCENT)
        styled_textbox(
            slide, left + Inches(0.12), top + Inches(0.07), w - Inches(0.24), Inches(0.28),
            head, size=11, bold=True, color=WHITE,
        )
        styled_textbox(
            slide, left + Inches(0.12), top + Inches(0.48), w - Inches(0.24), card_h - Inches(0.56),
            body, size=9, color=DARK,
        )


HOME_FEATURES = [
    (
        "Explore NDCs",
        "Planners · sector ministries · focal points",
        "Browse sector targets, Climate TRACE progress, mitigation options, and finance links in one cockpit.\n\n"
        "Useful when you need to see which sectors are on or off track before BTR reporting.",
    ),
    (
        "Data Ingestion",
        "MRV teams · GIS · ministry statistics",
        "Quick-scan unknown files or map columns to indicators and publish trusted observations.\n\n"
        "Useful when ministry spreadsheets must sit beside satellite data — triage without committing bad rows.",
    ),
    (
        "AI Predictions",
        "Strategy planners · finance · gap screening",
        "Forecast 2030 trajectories from observed emissions history; gap cards flag emerging shortfalls.\n\n"
        "Useful for prioritisation conversations — spot problems early, not at the 2030 deadline.",
    ),
    (
        "Policy documents",
        "Policy analysts · legal · delegation prep",
        "National corpus, CPR passage search on key laws, and intervention pathway diagrams.\n\n"
        "Useful to find evidence in NDC, NDPIV, and regulations — trace interventions to intended outcomes.",
    ),
    (
        "Policy Impact",
        "Socio-economic planners · adaptation teams",
        "Model policy choices with international case-study analogies and TEF intervention picker.\n\n"
        "Useful to stress-test options before committing resources — deep-link to finance screening.",
    ),
    (
        "My Work",
        "Field officers · delivery units · validators",
        "Log delivery activities, manage approvals queue, and track pending verifications.\n\n"
        "Useful to keep evidence organized before it counts toward national reporting.",
    ),
]


def add_six_feature_cards(slide, features, top=BODY_TOP):
    lefts, w, _ = col_widths(3)
    row_gap = Inches(0.16)
    avail_h = BODY_BOTTOM - top - row_gap
    card_h = avail_h / 2
    accents = [ACCENT, AMBER, RGBColor(155, 89, 182), MUTED, OCEAN, RGBColor(230, 126, 34)]
    for i, (title, audience, body) in enumerate(features):
        col = i % 3
        row = i // 3
        left = lefts[col]
        y = top + row * (card_h + row_gap)
        accent = accents[i % len(accents)]
        add_rect(slide, left, y, w, card_h, WHITE, PALE, radius=True)
        add_rect(slide, left, y, w, Inches(0.32), accent)
        styled_textbox(
            slide, left + Inches(0.1), y + Inches(0.05), w - Inches(0.2), Inches(0.24),
            title, size=10, bold=True, color=WHITE,
        )
        styled_textbox(
            slide, left + Inches(0.1), y + Inches(0.36), w - Inches(0.2), Inches(0.22),
            audience, size=7, bold=True, color=accent,
        )
        styled_textbox(
            slide, left + Inches(0.1), y + Inches(0.58), w - Inches(0.2), card_h - Inches(0.66),
            body, size=7.5, color=DARK,
        )


def slide_home_features(prs, num):
    s = content_slide(
        prs,
        "What you can do here",
        "Six modules for planners, MRV teams and partners — Climate TRACE + Uganda's Updated NDC",
        num,
    )
    add_rect(s, MARGIN_L, BODY_TOP, CONTENT_W, BODY_BOTTOM - BODY_TOP, LIGHT)
    add_six_feature_cards(s, HOME_FEATURES, top=BODY_TOP + Inches(0.1))
    styled_textbox(
        s, MARGIN_L + Inches(0.12), BODY_BOTTOM - Inches(0.42), CONTENT_W - Inches(0.24), Inches(0.32),
        "Each tile maps to a live route in the app — same home screen users see at " + LIVE_URL.replace("https://", ""),
        size=8, color=MUTED, align=PP_ALIGN.CENTER,
    )
    return s


def add_challenge_rows(slide, rows, top=BODY_TOP):
    row_h = Inches(0.58)
    y = top
    for title, fix in rows:
        add_rect(slide, MARGIN_L, y, CONTENT_W, row_h, WHITE, PALE)
        styled_textbox(
            slide, MARGIN_L + Inches(0.1), y + Inches(0.1), Inches(2.4), row_h - Inches(0.12),
            title, size=10, bold=True,
        )
        styled_textbox(
            slide, MARGIN_L + Inches(2.55), y + Inches(0.1), Inches(6.3), row_h - Inches(0.12),
            fix, size=9, color=MUTED,
        )
        y += row_h + Inches(0.06)


def add_flow_node(slide, left, top, width, height, title, subtitle, *, fill=WHITE, accent=ACCENT, title_color=DARK):
    add_rect(slide, left, top, width, height, fill, PALE, radius=True)
    add_rect(slide, left, top, width, Inches(0.06), accent)
    styled_textbox(
        slide, left + Inches(0.08), top + Inches(0.14), width - Inches(0.16), Inches(0.32),
        title, size=10, bold=True, color=title_color, align=PP_ALIGN.CENTER,
    )
    if subtitle:
        styled_textbox(
            slide, left + Inches(0.08), top + Inches(0.44), width - Inches(0.16), height - Inches(0.52),
            subtitle, size=8, color=MUTED, align=PP_ALIGN.CENTER,
        )


def add_h_arrow(slide, left, top, width):
    arrow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, left, top, width, Inches(0.22),
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ACCENT_LIGHT
    arrow.line.fill.background()


def add_v_arrow(slide, left, top, height):
    arrow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.DOWN_ARROW, left, top, Inches(0.22), height,
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ACCENT_LIGHT
    arrow.line.fill.background()


def add_connector_line(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = ACCENT_LIGHT
    conn.line.width = Pt(1.5)


def add_pipeline_row(slide, nodes, top, *, node_w=Inches(1.45), node_h=Inches(0.95), gap=Inches(0.28)):
    """Horizontal pipeline: list of (title, subtitle)."""
    arrow_w = gap
    total = len(nodes) * node_w + (len(nodes) - 1) * arrow_w
    left = MARGIN_L + (CONTENT_W - total) / 2
    x = left
    for i, (title, sub) in enumerate(nodes):
        add_flow_node(slide, x, top, node_w, node_h, title, sub)
        if i < len(nodes) - 1:
            add_h_arrow(slide, x + node_w, top + node_h / 2 - Inches(0.11), arrow_w)
        x += node_w + arrow_w


def add_screenshot(slide, filename, caption, num, *, subtitle=None, width=None, left=None):
    add_header(slide, "Live product", subtitle or caption)
    img_path = SCREENSHOTS / filename
    body_h = BODY_BOTTOM - BODY_TOP
    w = width or CONTENT_W
    l = left or MARGIN_L
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), l, BODY_TOP, width=w)
        pic = slide.shapes[-1]
        scale = min(w / pic.width, body_h / pic.height)
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)
        pic.left = int(l + (w - pic.width) / 2)
        pic.top = int(BODY_TOP + (body_h - pic.height) / 2)
    else:
        add_rect(slide, l, BODY_TOP, w, body_h, LIGHT, ACCENT)
        styled_textbox(
            slide, l, BODY_TOP + body_h / 2, w, Inches(0.3),
            f"[ Screenshot: {caption} ]", size=11, color=MUTED, align=PP_ALIGN.CENTER,
        )
    add_footer(slide, num)


def slide_problem_flow(prs, num):
    s = content_slide(prs, "The problem", "From static NDCs to operational delivery", num)
    add_rect(s, MARGIN_L, BODY_TOP, CONTENT_W, BODY_BOTTOM - BODY_TOP, LIGHT)
    styled_textbox(
        s, MARGIN_L + Inches(0.15), BODY_TOP + Inches(0.12), CONTENT_W - Inches(0.3), Inches(0.35),
        "Today: ministries assemble answers manually", size=11, bold=True, color=DARK,
    )
    nodes = [
        ("NDC PDFs", "Targets & pledges"),
        ("Spreadsheets", "Sector data silos"),
        ("Email / MRV", "Fragmented evidence"),
        ("Weeks of work", "Before any briefing"),
    ]
    add_pipeline_row(s, nodes, BODY_TOP + Inches(0.55), node_w=Inches(1.85), node_h=Inches(0.88), gap=Inches(0.2))

    mid_y = BODY_TOP + Inches(1.65)
    add_v_arrow(s, MARGIN_L + CONTENT_W / 2 - Inches(0.11), mid_y, Inches(0.35))
    styled_textbox(
        s, MARGIN_L, mid_y + Inches(0.42), CONTENT_W, Inches(0.3),
        "Five operational questions — unanswered in real time", size=10, bold=True, color=ACCENT, align=PP_ALIGN.CENTER,
    )

    questions = [
        "Where are we today?",
        "What moves the needle?",
        "Policy impact?",
        "Gaps vs target?",
        "What to prioritise?",
    ]
    qw = Inches(1.55)
    qgap = Inches(0.12)
    qtotal = len(questions) * qw + (len(questions) - 1) * qgap
    qx = MARGIN_L + (CONTENT_W - qtotal) / 2
    qy = mid_y + Inches(0.78)
    for q in questions:
        add_rect(s, qx, qy, qw, Inches(0.62), WHITE, ACCENT, radius=True)
        styled_textbox(s, qx + Inches(0.06), qy + Inches(0.18), qw - Inches(0.12), Inches(0.35),
                       q, size=8, color=DARK, align=PP_ALIGN.CENTER)
        qx += qw + qgap

    styled_textbox(
        s, MARGIN_L + Inches(0.15), BODY_BOTTOM - Inches(0.55), CONTENT_W - Inches(0.3), Inches(0.45),
        "Our answer: link Uganda's NDC layer to live observed emissions via the Climate TRACE API — "
        "not another static report.",
        size=9, color=MUTED,
    )
    return s


def slide_cockpit_architecture(prs, num):
    s = content_slide(prs, "Cockpit architecture", "Five linked columns — one decision workspace", num)
    add_rect(s, MARGIN_L, BODY_TOP, CONTENT_W, Inches(3.55), LIGHT)
    cols = [
        ("1 · Targets", "NDC pledges\nSector filters"),
        ("2 · Activities", "Measures & owners\nImplementation catalog"),
        ("3 · Observed", "Climate TRACE API\nLive MtCO₂e"),
        ("4 · Progress", "On-track scoring\nGap vs 2030"),
        ("5 · Options", "Mitigation paths\nFinance scenarios"),
    ]
    cw = Inches(1.55)
    cg = Inches(0.12)
    cx = MARGIN_L + Inches(0.35)
    cy = BODY_TOP + Inches(0.35)
    for title, sub in cols:
        fill = RGBColor(235, 245, 238) if "Observed" in title else WHITE
        accent = OCEAN if "Observed" in title else ACCENT
        add_flow_node(s, cx, cy, cw, Inches(1.35), title, sub, fill=fill, accent=accent)
        if title != cols[-1][0]:
            add_h_arrow(s, cx + cw, cy + Inches(0.56), cg)
        cx += cw + cg

    # Vertical feeds below
    feeds = [
        (MARGIN_L + Inches(0.2), "Uganda NDC catalog", "Curated targets & strategies"),
        (MARGIN_L + Inches(2.45), "Climate TRACE v7", "api.climatetrace.org"),
        (MARGIN_L + Inches(4.7), "3D GIS map tab", "Satellite terrain\nEmission bubbles"),
        (MARGIN_L + Inches(6.95), "NDC AI chatbot", "Grounded briefing\nOpenAI + citations"),
    ]
    fy = BODY_TOP + Inches(2.05)
    fw = Inches(2.05)
    for fx, ft, fs in feeds:
        accent = OCEAN if "TRACE" in ft else (AMBER if "AI" in ft else ACCENT_LIGHT)
        add_flow_node(s, fx, fy, fw, Inches(0.88), ft, fs, fill=WHITE, accent=accent)

    for fx in [MARGIN_L + Inches(1.2), MARGIN_L + Inches(3.45), MARGIN_L + Inches(5.7)]:
        add_connector_line(s, fx, fy, fx, cy + Inches(1.35))

    styled_textbox(
        s, MARGIN_L + Inches(0.15), BODY_TOP + Inches(3.15), CONTENT_W - Inches(0.3), Inches(0.55),
        "Observed column + Emissions Map + NDC AI share one Climate TRACE pipeline — "
        "charts, 3D GIS bubbles, and chat answers cite the same live API data.",
        size=9, color=DARK,
    )
    return s


def slide_climate_trace_pipeline(prs, num):
    s = content_slide(prs, "Climate TRACE data pipeline", "Live observed emissions — API v7 (CC BY 4.0)", num)
    add_rect(s, MARGIN_L, BODY_TOP, CONTENT_W, BODY_BOTTOM - BODY_TOP, LIGHT)

    # Source layer
    styled_textbox(s, MARGIN_L + Inches(0.12), BODY_TOP + Inches(0.1), Inches(2), Inches(0.25),
                   "EMISSIONS SOURCES", size=8, bold=True, color=MUTED)
    sources = [
        ("Satellite EO", "Land & power"),
        ("Facility data", "Power & industry"),
        ("Transport", "Road & aviation"),
        ("AFOLU models", "Land use"),
    ]
    sx = MARGIN_L + Inches(0.12)
    sy = BODY_TOP + Inches(0.38)
    for title, sub in sources:
        add_flow_node(s, sx, sy, Inches(1.95), Inches(0.72), title, sub, fill=RGBColor(230, 242, 248), accent=OCEAN)
        sy += Inches(0.82)

    # API box (centre hero)
    api_l = MARGIN_L + Inches(2.55)
    api_t = BODY_TOP + Inches(0.55)
    api_w = Inches(3.35)
    api_h = Inches(2.65)
    add_rect(s, api_l, api_t, api_w, api_h, WHITE, OCEAN, radius=True)
    add_rect(s, api_l, api_t, api_w, Inches(0.42), OCEAN)
    styled_textbox(s, api_l + Inches(0.12), api_t + Inches(0.08), api_w - Inches(0.24), Inches(0.3),
                   "Climate TRACE API v7", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    endpoints = [
        "GET /v7/sources/emissions — sector timeseries",
        "GET /v7/sources — asset-level map points",
        "GET /v7/rankings/countries — national totals",
        "GADM2 districts — Uganda disaggregation",
    ]
    styled_textbox(
        s, api_l + Inches(0.15), api_t + Inches(0.52), api_w - Inches(0.3), api_h - Inches(0.6),
        "\n".join(f"•  {e}" for e in endpoints),
        size=8, color=DARK,
    )
    styled_textbox(s, api_l + Inches(0.15), api_t + api_h - Inches(0.38), api_w - Inches(0.3), Inches(0.3),
                   CLIMATE_TRACE_API, size=7, color=OCEAN, align=PP_ALIGN.CENTER)

    # Arrows sources → API
    for i in range(4):
        y = BODY_TOP + Inches(0.74) + i * Inches(0.82)
        add_connector_line(s, MARGIN_L + Inches(2.07), y, api_l, api_t + Inches(1.3))

    # Our backend
    back_l = MARGIN_L + Inches(6.35)
    back_nodes = [
        ("Express API", "Node.js backend\nCache & provenance"),
        ("Postgres / seed", "MtCO₂e store\n/api/v1/provenance"),
        ("React cockpit", "Dashboard · 3D GIS\nNDC AI chatbot"),
    ]
    by = BODY_TOP + Inches(0.38)
    for title, sub in back_nodes:
        add_flow_node(s, back_l, by, Inches(2.45), Inches(0.88), title, sub)
        add_h_arrow(s, api_l + api_w, by + Inches(0.33), Inches(0.42))
        by += Inches(1.02)

    styled_textbox(
        s, MARGIN_L + Inches(0.12), BODY_BOTTOM - Inches(0.62), CONTENT_W - Inches(0.24), Inches(0.5),
        "In the app: every observed chart shows Climate TRACE provenance — click a bar to trace the source. "
        "Not a national inventory replacement; open observed layer for briefing.",
        size=8, color=MUTED,
    )
    return s


def slide_ai_and_3d_gis(prs, num):
    """NDC AI chatbot + Emissions Map 3D GIS representation tab."""
    s = content_slide(prs, "NDC AI & 3D GIS", "Intelligence and spatial views on Climate TRACE data", num)
    add_rect(s, MARGIN_L, BODY_TOP, CONTENT_W, BODY_BOTTOM - BODY_TOP, LIGHT)

    half = CONTENT_W / 2 - Inches(0.1)
    left = MARGIN_L
    right = MARGIN_L + half + Inches(0.2)
    panel_h = Inches(3.55)
    panel_t = BODY_TOP + Inches(0.12)

    # ── Left: 3D GIS Representation (Emissions Map tab) ──
    add_rect(s, left, panel_t, half, panel_h, WHITE, OCEAN, radius=True)
    add_rect(s, left, panel_t, half, Inches(0.36), OCEAN)
    styled_textbox(s, left + Inches(0.12), panel_t + Inches(0.06), half - Inches(0.24), Inches(0.28),
                   "3D GIS representation — Emissions Map", size=11, bold=True, color=WHITE)
    gis_flow = [
        ("Climate TRACE\n/v7/sources", "Asset geo-points\nper sector"),
        ("MapLibre GL", "Satellite imagery\nTerrain DEM 3D"),
        ("GIS bubbles", "Size ∝ MtCO₂e\nDistrict overlay"),
    ]
    gx = left + Inches(0.15)
    gy = panel_t + Inches(0.48)
    gw = Inches(1.35)
    for title, sub in gis_flow:
        add_flow_node(s, gx, gy, gw, Inches(0.78), title, sub, fill=RGBColor(230, 242, 248), accent=OCEAN)
        if title != gis_flow[-1][0]:
            add_h_arrow(s, gx + gw, gy + Inches(0.28), Inches(0.1))
        gx += gw + Inches(0.1)

    img_path = SCREENSHOTS / "map-uganda.png"
    img_l = left + Inches(0.15)
    img_t = panel_t + Inches(1.38)
    img_w = half - Inches(0.3)
    img_h = panel_h - Inches(1.5)
    add_rect(s, img_l, img_t, img_w, img_h, RGBColor(12, 16, 24), PALE, radius=True)
    if img_path.exists():
        s.shapes.add_picture(str(img_path), img_l + Inches(0.04), img_t + Inches(0.04), width=img_w - Inches(0.08))
        pic = s.shapes[-1]
        scale = min((img_w - Inches(0.08)) / pic.width, (img_h - Inches(0.08)) / pic.height)
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)
        pic.left = int(img_l + (img_w - pic.width) / 2)
        pic.top = int(img_t + (img_h - pic.height) / 2)
    else:
        styled_textbox(s, img_l, img_t + img_h / 2, img_w, Inches(0.25),
                       "[ Emissions Map — 3D GIS ]", size=9, color=MUTED, align=PP_ALIGN.CENTER)

    # ── Right: NDC AI chatbot ──
    add_rect(s, right, panel_t, half, panel_h, WHITE, ACCENT, radius=True)
    add_rect(s, right, panel_t, half, Inches(0.36), ACCENT)
    styled_textbox(s, right + Inches(0.12), panel_t + Inches(0.06), half - Inches(0.24), Inches(0.28),
                   "NDC AI chatbot — dashboard sidebar", size=11, bold=True, color=WHITE)

    ai_flow = [
        ("User question", "Plain language\n'Transport on track?'"),
        ("Context pack", "Sector · target\nTRACE progress"),
        ("OpenAI API", "Grounded answer\nWith citations"),
    ]
    ax = right + Inches(0.15)
    ay = panel_t + Inches(0.48)
    for title, sub in ai_flow:
        add_flow_node(s, ax, ay, gw, Inches(0.78), title, sub, fill=RGBColor(235, 245, 238), accent=ACCENT)
        if title != ai_flow[-1][0]:
            add_h_arrow(s, ax + gw, ay + Inches(0.28), Inches(0.1))
        ax += gw + Inches(0.1)

    ai_bullets = [
        "Opens from Progress column → NDC AI button",
        "Tab key fills an example question",
        "Answers cite Climate TRACE + NDC panels",
        "Senior Decision-Maker briefing mode",
        "Not a generic chatbot — country-specific context",
    ]
    add_rect(s, right + Inches(0.15), panel_t + Inches(1.38), half - Inches(0.3), img_h, RGBColor(248, 250, 249), PALE, radius=True)
    styled_textbox(
        s, right + Inches(0.28), panel_t + Inches(1.52), half - Inches(0.55), img_h - Inches(0.2),
        "\n".join(f"▸  {b}" for b in ai_bullets),
        size=8, color=DARK,
    )

    styled_textbox(
        s, MARGIN_L + Inches(0.12), BODY_BOTTOM - Inches(0.48), CONTENT_W - Inches(0.24), Inches(0.38),
        "Both modules sit on the same Climate TRACE API v7 feed — spatial GIS for where emissions occur, "
        "NDC AI for why it matters against Uganda's pledges.",
        size=8, color=MUTED, align=PP_ALIGN.CENTER,
    )
    return s


def slide_satellite_to_decision(prs, num):
    s = content_slide(prs, "From satellite to decision", "How observed data reaches the dashboard", num)
    steps = [
        ("Climate TRACE\ningests signals", "Monthly updates\n100+ countries"),
        ("API v7 serves\nUganda sectors", "Transport · Energy\nAFOLU · Waste…"),
        ("App reconciles\nvs NDC targets", "Progress % · gaps\nDistrict views"),
        ("Decision-maker\nbriefing", "Dashboard · Map\nNDC AI · Finance"),
    ]
    add_pipeline_row(s, steps, BODY_TOP + Inches(0.25), node_w=Inches(1.95), node_h=Inches(1.05), gap=Inches(0.18))

    # Screenshot inset
    img_path = SCREENSHOTS / "map-uganda.png"
    inset_l = MARGIN_L + Inches(0.15)
    inset_t = BODY_TOP + Inches(1.55)
    inset_w = Inches(4.2)
    inset_h = BODY_BOTTOM - inset_t - Inches(0.15)
    add_rect(s, inset_l, inset_t, inset_w, inset_h, WHITE, PALE, radius=True)
    if img_path.exists():
        s.shapes.add_picture(str(img_path), inset_l + Inches(0.08), inset_t + Inches(0.08),
                             width=inset_w - Inches(0.16))
        pic = s.shapes[-1]
        scale = min((inset_w - Inches(0.16)) / pic.width, (inset_h - Inches(0.16)) / pic.height)
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)
        pic.left = int(inset_l + (inset_w - pic.width) / 2)
        pic.top = int(inset_t + (inset_h - pic.height) / 2)
    else:
        styled_textbox(s, inset_l, inset_t + inset_h / 2, inset_w, Inches(0.3),
                       "[ Emissions map ]", size=10, color=MUTED, align=PP_ALIGN.CENTER)

    legend_l = MARGIN_L + Inches(4.55)
    legend_t = inset_t
    legend_w = CONTENT_W - Inches(4.7)
    add_rect(s, legend_l, legend_t, legend_w, inset_h, WHITE, ACCENT, radius=True)
    styled_textbox(s, legend_l + Inches(0.15), legend_t + Inches(0.15), legend_w - Inches(0.3), Inches(0.3),
                   "What users see", size=11, bold=True, color=ACCENT)
    bullets = [
        "Live Climate TRACE banner on dashboard",
        "Sector timeseries charts (MtCO₂e)",
        "3D GIS map — satellite terrain + emission bubbles",
        "NDC AI chatbot — grounded briefing with citations",
        "Progress ring vs 2030 NDC pledge",
    ]
    styled_textbox(s, legend_l + Inches(0.15), legend_t + Inches(0.5), legend_w - Inches(0.3), inset_h - Inches(0.6),
                   "\n\n".join(f"▸  {b}" for b in bullets), size=9, color=DARK)
    return s


def slide_demo_journey(prs, num):
    s = content_slide(prs, "5-minute live walkthrough", "Suggested flow — Uganda · Transport sector", num)
    add_rect(s, MARGIN_L, BODY_TOP, CONTENT_W, BODY_BOTTOM - BODY_TOP, LIGHT)
    journey = [
        ("0:00", "Explore NDCs", "Dashboard · Transport\nClimate TRACE progress"),
        ("0:45", "Data Ingestion", "Quick scan file\nMapped publish"),
        ("1:15", "AI Predictions", "2030 trajectory\nGap cards"),
        ("1:45", "Policy docs", "CPR passages\nIntervention path"),
        ("2:15", "Policy Impact", "KCI analogies\nFinance link"),
        ("2:45", "My Work", "Activities queue\nVerifications"),
    ]
    jw = Inches(1.22)
    jh = Inches(1.05)
    jgap = Inches(0.1)
    jx = MARGIN_L + Inches(0.12)
    jy = BODY_TOP + Inches(0.28)
    for time, title, sub in journey:
        add_rect(s, jx, jy, jw, jh, WHITE, ACCENT, radius=True)
        add_rect(s, jx, jy, jw, Inches(0.28), ACCENT)
        styled_textbox(s, jx, jy + Inches(0.04), jw, Inches(0.22), time, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        styled_textbox(s, jx + Inches(0.06), jy + Inches(0.34), jw - Inches(0.12), Inches(0.3), title, size=9, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        styled_textbox(s, jx + Inches(0.06), jy + Inches(0.62), jw - Inches(0.12), Inches(0.45), sub, size=7, color=MUTED, align=PP_ALIGN.CENTER)
        if title != journey[-1][1]:
            add_h_arrow(s, jx + jw, jy + jh / 2 - Inches(0.11), jgap)
        jx += jw + jgap

    styled_textbox(
        s, MARGIN_L + Inches(0.2), BODY_TOP + Inches(1.55), CONTENT_W - Inches(0.4), Inches(0.35),
        "Punchline: all six home modules — from observed emissions to policy evidence and delivery workflow.",
        size=10, bold=True, color=ACCENT, align=PP_ALIGN.CENTER,
    )

    # Bottom row: extras + live URL
    future = [
        ("Also in app", "3D GIS map · NDC AI chatbot\nFinance & MAC screening"),
        ("Challenges", "Mixed sources · API warm-up\nNot UNFCCC submission"),
        ("Live URL", LIVE_URL.replace("https://", "")),
    ]
    fx = MARGIN_L + Inches(0.12)
    fy = BODY_TOP + Inches(2.05)
    fw = Inches(2.75)
    for title, sub in future:
        add_flow_node(s, fx, fy, fw, Inches(0.95), title, sub)
        fx += fw + Inches(0.2)
    return s


def slide_system_architecture(prs, num):
    s = content_slide(prs, "System architecture", "Browser → API → live data — same code local and on Vercel", num)
    add_rect(s, MARGIN_L, BODY_TOP, CONTENT_W, BODY_BOTTOM - BODY_TOP, LIGHT)

    band_l = MARGIN_L + Inches(0.15)
    band_w = CONTENT_W - Inches(0.3)
    label_w = Inches(1.5)
    node_area_l = band_l + label_w + Inches(0.12)
    node_area_w = band_w - label_w - Inches(0.12)
    arrow_x = node_area_l + node_area_w / 2 - Inches(0.11)

    def band(top, h, label, sublabel, accent, nodes, node_fill=WHITE):
        add_rect(s, band_l, top, band_w, h, WHITE, PALE)
        add_rect(s, band_l, top, label_w, h, accent)
        styled_textbox(s, band_l + Inches(0.1), top + Inches(0.12), label_w - Inches(0.2), Inches(0.3),
                       label, size=10, bold=True, color=WHITE)
        styled_textbox(s, band_l + Inches(0.1), top + h - Inches(0.4), label_w - Inches(0.2), Inches(0.32),
                       sublabel, size=7, color=PALE)
        n = len(nodes)
        gap = Inches(0.1)
        nw = (node_area_w - gap * (n - 1)) / n
        nx = node_area_l
        ny = top + Inches(0.16)
        for title, sub in nodes:
            add_rect(s, nx, ny, nw, h - Inches(0.32), node_fill, PALE, radius=True)
            styled_textbox(s, nx + Inches(0.05), ny + Inches(0.08), nw - Inches(0.1), Inches(0.28),
                           title, size=9, bold=True, color=DARK, align=PP_ALIGN.CENTER)
            if sub:
                styled_textbox(s, nx + Inches(0.05), ny + Inches(0.34), nw - Inches(0.1), h - Inches(0.66),
                               sub, size=7, color=MUTED, align=PP_ALIGN.CENTER)
            nx += nw + gap

    top1 = BODY_TOP + Inches(0.14)
    h = Inches(1.0)
    band(top1, h, "Presentation", "React SPA", ACCENT, [
        ("Dashboard", "Targets · observed · progress"),
        ("Emissions Map", "3D GIS bubbles"),
        ("Policy / Finance", "Pathways · MAC"),
        ("NDC AI", "Grounded chat"),
    ])
    add_v_arrow(s, arrow_x, top1 + h + Inches(0.02), Inches(0.26))
    styled_textbox(s, arrow_x + Inches(0.22), top1 + h + Inches(0.0), Inches(3.2), Inches(0.24),
                   "same-origin /api/v1 · React Query cache", size=7, color=MUTED)

    top2 = top1 + h + Inches(0.32)
    band(top2, h, "API layer", "Express · Node", ACCENT_LIGHT, [
        ("emissions", "TRACE aggregation"),
        ("documents", "CPR corpus"),
        ("policyImpact", "KCI matching"),
        ("ingest", "Mapped import"),
        ("dashboardAi", "OpenAI proxy"),
    ])
    add_v_arrow(s, arrow_x, top2 + h + Inches(0.02), Inches(0.26))

    top3 = top2 + h + Inches(0.32)
    band(top3, h, "Data sources", "live + versioned", OCEAN, [
        ("Climate TRACE v7", "Live emissions API"),
        ("NDC catalog", "Repo-versioned"),
        ("Postgres", "Optional ingest"),
        ("OpenAI", "gpt-4o-mini"),
    ], node_fill=RGBColor(235, 245, 238))
    return s


def slide_three_layers(prs, num):
    s = content_slide(prs, "Three layers of truth", "What is official, what is observed, what is indicative", num)
    add_rect(s, MARGIN_L, BODY_TOP, CONTENT_W, BODY_BOTTOM - BODY_TOP, LIGHT)
    rows = [
        (ACCENT, "1 · Official pledges",
         "Uganda Updated NDC 2022 targets + curated activity catalogue — versioned in the repo.", "BUNDLED"),
        (OCEAN, "2 · Observed emissions",
         "Climate TRACE API v7 — live MtCO₂e on Dashboard & 3D GIS Map; national 2015+, district 2021+.", "LIVE"),
        (AMBER, "3 · Evidence & screening",
         "Policy docs (CPR export), indicative finance (MAC), Policy Impact (KCI analogies), pathway diagrams.", "INDICATIVE"),
        (ACCENT_LIGHT, "4 · Ministry uploads",
         "Mapped ingest observations on indicator targets — provenance badge shown when Postgres is configured.", "OPTIONAL"),
    ]
    rh = Inches(0.78)
    gap = Inches(0.14)
    y = BODY_TOP + Inches(0.16)
    badge_w = Inches(1.1)
    for color, title, desc, tag in rows:
        add_rect(s, MARGIN_L + Inches(0.15), y, CONTENT_W - Inches(0.3), rh, WHITE, PALE, radius=True)
        add_rect(s, MARGIN_L + Inches(0.15), y, Inches(0.12), rh, color)
        styled_textbox(s, MARGIN_L + Inches(0.4), y + Inches(0.1), Inches(2.7), Inches(0.3),
                       title, size=11, bold=True, color=color)
        styled_textbox(s, MARGIN_L + Inches(0.4), y + Inches(0.4), Inches(6.4), Inches(0.34),
                       desc, size=8.5, color=DARK)
        badge_l = MARGIN_L + CONTENT_W - Inches(0.25) - badge_w
        add_rect(s, badge_l, y + Inches(0.22), badge_w, Inches(0.34), color, radius=True)
        styled_textbox(s, badge_l, y + Inches(0.27), badge_w, Inches(0.24),
                       tag, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        y += rh + gap
    styled_textbox(s, MARGIN_L + Inches(0.15), y + Inches(0.02), CONTENT_W - Inches(0.3), Inches(0.3),
                   "Intended outcomes (targets, pathways) stay visually distinct from measured outcomes (Climate TRACE).",
                   size=8, color=MUTED, align=PP_ALIGN.CENTER)
    return s


def slide_dashboard(prs, num):
    s = content_slide(prs, "Live product — NDC Dashboard", "Observed emissions linked to NDC targets", num)
    body_h = BODY_BOTTOM - BODY_TOP
    img_w = Inches(6.05)
    img_l = MARGIN_L
    add_rect(s, img_l, BODY_TOP, img_w, body_h, WHITE, PALE, radius=True)
    img_path = SCREENSHOTS / "dashboard-transport.png"
    if img_path.exists():
        s.shapes.add_picture(str(img_path), img_l + Inches(0.08), BODY_TOP + Inches(0.08), width=img_w - Inches(0.16))
        pic = s.shapes[-1]
        scale = min((img_w - Inches(0.16)) / pic.width, (body_h - Inches(0.16)) / pic.height)
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)
        pic.left = int(img_l + (img_w - pic.width) / 2)
        pic.top = int(BODY_TOP + (body_h - pic.height) / 2)
    else:
        styled_textbox(s, img_l, BODY_TOP + body_h / 2, img_w, Inches(0.3),
                       "[ Dashboard — Transport ]", size=11, color=MUTED, align=PP_ALIGN.CENTER)
    px = img_l + img_w + Inches(0.2)
    pw = MARGIN_L + CONTENT_W - px
    add_rect(s, px, BODY_TOP, pw, body_h, WHITE, ACCENT, radius=True)
    add_rect(s, px, BODY_TOP, pw, Inches(0.36), ACCENT)
    styled_textbox(s, px + Inches(0.12), BODY_TOP + Inches(0.06), pw - Inches(0.24), Inches(0.28),
                   "On screen", size=11, bold=True, color=WHITE)
    points = [
        ("Targets", "Uganda NDC 2022 pledges by sector"),
        ("Observed", "Live Climate TRACE MtCO₂e timeseries"),
        ("Progress", "On-track status vs 2030 ceiling"),
        ("NDC AI", "Ask a grounded question in plain language"),
        ("Provenance", "Click any bar to trace the source"),
    ]
    py = BODY_TOP + Inches(0.5)
    for head, sub in points:
        styled_textbox(s, px + Inches(0.14), py, pw - Inches(0.28), Inches(0.24),
                       head, size=9.5, bold=True, color=ACCENT)
        styled_textbox(s, px + Inches(0.14), py + Inches(0.22), pw - Inches(0.28), Inches(0.4),
                       sub, size=8, color=DARK)
        py += Inches(0.62)
    return s


def slide_policy_finance(prs, num):
    s = content_slide(prs, "Policy & finance modules", "From observed gaps to investment scenarios", num)
    half_w = CONTENT_W / 2 - Inches(0.12)
    mid = MARGIN_L + half_w + Inches(0.24)
    body_h = BODY_BOTTOM - BODY_TOP
    card_h = body_h - Inches(0.95)
    cards = [
        (SCREENSHOTS / "documents-pathway.png", MARGIN_L, "Policy pathway",
         "Interventions → outcomes, traced to NDC targets (TEF logic model)."),
        (SCREENSHOTS / "climate-finance.png", mid, "Climate finance",
         "Indicative MAC — abatement cost vs potential; fund + MCF hints."),
    ]
    for path, left, label, desc in cards:
        add_rect(s, left, BODY_TOP, half_w, card_h, WHITE, PALE, radius=True)
        if path.exists():
            s.shapes.add_picture(str(path), left + Inches(0.06), BODY_TOP + Inches(0.06), width=half_w - Inches(0.12))
            pic = s.shapes[-1]
            scale = min((half_w - Inches(0.12)) / pic.width, (card_h - Inches(0.5)) / pic.height)
            pic.width = int(pic.width * scale)
            pic.height = int(pic.height * scale)
            pic.left = int(left + (half_w - pic.width) / 2)
            pic.top = int(BODY_TOP + Inches(0.06))
        else:
            styled_textbox(s, left, BODY_TOP + card_h / 2, half_w, Inches(0.3),
                           f"[ {label} ]", size=10, color=MUTED, align=PP_ALIGN.CENTER)
        styled_textbox(s, left + Inches(0.12), BODY_TOP + card_h - Inches(0.34), half_w - Inches(0.24), Inches(0.3),
                       desc, size=8, color=MUTED)
    add_pipeline_row(s, [
        ("Climate TRACE", "Observed gaps"),
        ("Policy docs", "Intervention map"),
        ("Finance MAC", "Abatement cost"),
    ], BODY_TOP + card_h + Inches(0.12), node_w=Inches(2.55), node_h=Inches(0.62), gap=Inches(0.35))
    return s


def build():
    prs = new_prs()

    # 01 Title
    s = blank_slide(prs)
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, DARK)
    add_rect(s, Inches(0), Inches(4.05), SLIDE_W, Inches(1.575), ACCENT)
    styled_textbox(s, MARGIN_L, Inches(1.35), CONTENT_W, Inches(0.7),
                   "Uganda NDC Data Explorer", size=32, bold=True, color=WHITE)
    styled_textbox(s, MARGIN_L, Inches(2.05), CONTENT_W, Inches(0.45),
                   "Decision-support cockpit · powered by Climate TRACE API v7", size=14, color=SUBTITLE)
    styled_textbox(s, MARGIN_L, Inches(2.65), CONTENT_W, Inches(0.55),
                   "Climate TRACE API  →  3D GIS map  →  NDC AI chatbot  →  policy decisions",
                   size=11, color=SUBTITLE)
    styled_textbox(s, MARGIN_L, Inches(3.25), CONTENT_W, Inches(0.4),
                   "Stack: React · Express (Node) · Climate TRACE API v7 · OpenAI · MapLibre GL · Vercel",
                   size=10, color=SUBTITLE)
    styled_textbox(s, MARGIN_L, Inches(4.25), CONTENT_W, Inches(0.9),
                   "Date: 17 June 2026\nPresented by: Tseten Sherpa",
                   size=11, color=WHITE)

    # 02 Problem flow
    slide_problem_flow(prs, 2)

    # 03 Cockpit architecture
    slide_cockpit_architecture(prs, 3)

    # 04 Home — six feature modules (user value)
    slide_home_features(prs, 4)

    # 05 System architecture (tech stack diagram)
    slide_system_architecture(prs, 5)

    # 06 Climate TRACE pipeline (hero technical slide)
    slide_climate_trace_pipeline(prs, 6)

    # 07 Three layers of truth (data honesty diagram)
    slide_three_layers(prs, 7)

    # 08 Dashboard screenshot + annotation points
    slide_dashboard(prs, 8)

    # 09 NDC AI + 3D GIS (dedicated)
    slide_ai_and_3d_gis(prs, 9)

    # 10 End-to-end user value
    slide_satellite_to_decision(prs, 10)

    # 11 Policy & finance (screenshots + descriptions + mini flow)
    slide_policy_finance(prs, 11)

    # 12 Live walkthrough flowchart
    slide_demo_journey(prs, 12)

    # 13 Challenges
    s = content_slide(prs, "Challenges & mitigations", "Honest limits of a prototype", 13)
    add_challenge_rows(s, [
        ("Climate TRACE ≠ inventory", "Open observed layer; label provenance clearly"),
        ("API cold starts", "Pre-warm dashboard, 3D GIS map, and NDC AI before stage"),
        ("NDC AI availability", "Requires OPENAI_API_KEY on API server"),
        ("Mixed data sources", "Curated NDC catalog + live API panels"),
        ("Official reporting", "Briefing tool — not UNFCCC submission"),
    ])

    # 14 Future + close
    s = content_slide(prs, "Future roadmap", "Scaling beyond Uganda", 14)
    add_three_cards(s, [
        ("Policy paper database",
         "Searchable national policy corpus per country.\n\n"
         "Full-text search linked to NDC targets and Climate TRACE sectors."),
        ("Mobile application",
         "Field MRV capture and district briefings.\n\n"
         "Offline views when API is unavailable."),
        ("Multi-country scale",
         "Replicate cockpit for other NDCs.\n\n"
         "Shared Climate TRACE v7 pipeline; country strategy layers."),
    ], card_h=Inches(3.2))
    styled_textbox(
        s, MARGIN_L, Inches(4.72), CONTENT_W, Inches(0.38),
        f"Questions?  ·  {LIVE_URL}",
        size=9, color=MUTED, align=PP_ALIGN.CENTER,
    )

    prs.save(OUT)
    import shutil
    downloads = Path.home() / "Downloads" / OUT.name
    shutil.copy(OUT, downloads)
    print(f"Saved: {OUT}")
    print(f"Copied: {downloads}")


if __name__ == "__main__":
    build()
